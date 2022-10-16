
import pptx

   
def ppt2text(fname):
    prs = pptx.Presentation(fname)

    for i, sld in enumerate(prs.slides, start=1):

        print(f'-- Page {i} --')

        for shp in sld.shapes:

            if shp.has_text_frame:
                print (shp.text)

            if shp.has_table:
                tbl = shp.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0, row_count):                 
                    text=''
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        paragraphs = cell.text_frame.paragraphs 
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                text+=run.text
                            text+=', '
                    print (text)
            print ()
