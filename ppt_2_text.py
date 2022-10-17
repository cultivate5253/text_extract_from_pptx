
import os,sys
from tkinter import *
from tkinter import ttk
from tkinter import messagebox
from tkinter import filedialog

import pptx
from glob import glob

#for fname in glob ('*.pptx'):
#print ('File name: ', fname, '\n')
   
def ppt2text(fname):
    prs = pptx.Presentation(fname)

    for i, sld in enumerate(prs.slides, start=1):

        print(f'-- Page {i} --')

        for shp in sld.shapes:

            if shp.has_text_frame:
                print (shp.text)

            if shp.has_table:
                tbl = shp.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0, row_count):                 
                    text=''
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        paragraphs = cell.text_frame.paragraphs 
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                text+=run.text
                            text+=', '
                    print (text)
            print ()

# # フォルダ指定の関数
# def dirdialog_clicked():
#     iDir = os.path.dirname(os.path.abspath('__file__'))
#     iDirPath = filedialog.askdirectory(initialdir = iDir)
#     entry1.set(iDirPath)

# ファイル指定の関数
def filedialog_clicked():
    fTyp = [("", "*")]
    iFile =os.path.dirname(os.path.abspath('__file__'))
    iFilePath = filedialog.askopenfilename(filetype = fTyp, initialdir = iFile)
    entry2.set(iFilePath)

# 実行ボタン押下時の実行関数
def conductMain():
    text = ""

    # dirPath = entry1.get()
    filePath = entry2.get()
    # if dirPath:
    #     text += "フォルダパス：" + dirPath + "\n"
    if filePath:
        text += "ファイルパス：" + filePath

    if text:
       ppt2text(filePath)
    else:
        messagebox.showerror("error", "パスの指定がありません。")

if __name__ == "__main__":

    # rootの作成
    root = Tk()
    root.title("サンプル")

    # Frame1の作成
    frame1 = ttk.Frame(root, padding=10)
    frame1.grid(row=0, column=1, sticky=E)

    # # 「フォルダ参照」ラベルの作成
    # IDirLabel = ttk.Label(frame1, text="フォルダ参照＞＞", padding=(5, 2))
    # IDirLabel.pack(side=LEFT)

    # # 「フォルダ参照」エントリーの作成
    # entry1 = StringVar()
    # IDirEntry = ttk.Entry(frame1, textvariable=entry1, width=30)
    # IDirEntry.pack(side=LEFT)

    # # 「フォルダ参照」ボタンの作成
    # IDirButton = ttk.Button(frame1, text="参照", command=dirdialog_clicked)
    # IDirButton.pack(side=LEFT)

    # #Frame1の作成
    # frame1 = ttk.Frame(root, padding=10)
    # frame1.grid(row=2, column=1, sticky=E)

     #Frame2の作成
    frame2 = ttk.Frame(root, padding=10)
    frame2.grid(row=2, column=1, sticky=E)

    # 「ファイル参照」ラベルの作成
    IFileLabel = ttk.Label(frame2, text="ファイル参照＞＞", padding=(5, 2))
    IFileLabel.pack(side=LEFT)

    # 「ファイル参照」エントリーの作成
    entry2 = StringVar()
    IFileEntry = ttk.Entry(frame2, textvariable=entry2, width=30)
    IFileEntry.pack(side=LEFT)

    # 「ファイル参照」ボタンの作成
    IFileButton = ttk.Button(frame2, text="参照", command=filedialog_clicked)
    IFileButton.pack(side=LEFT)

    # Frame2の作成
    frame3 = ttk.Frame(root, padding=10)
    frame3.grid(row=5,column=1,sticky=W)

    # 実行ボタンの設置
    button1 = ttk.Button(frame3, text="実行", command=conductMain)
    button1.pack(fill = "x", padx=30, side = "left")

    # キャンセルボタンの設置
    button2 = ttk.Button(frame2, text=("閉じる"), command=quit)
    button2.pack(fill = "x", padx=30, side = "left")

    root.mainloop()

    input()
    